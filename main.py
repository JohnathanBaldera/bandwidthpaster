from tracemalloc import start
import openpyxl
from openpyxl.styles import PatternFill
from pptx import Presentation
from pptx.util import Inches
from tkinter import *
from tkinter import ttk
from tkinter import filedialog
import matplotlib.pyplot as plt


class BandwidthPaster:

    def __init__(self, root):
        root.title("Bandwidth Paster")

        mainframe = ttk.Frame(root, padding="3 4 12 12")
        mainframe.grid(column=0, row=0, sticky=(N, W, E, S))
        root.columnconfigure(0, weight=1)
        root.rowconfigure(0, weight=1)

        self.title = StringVar()
        title_entry = ttk.Entry(mainframe, textvariable=self.title)
        title_entry.grid(column=2, row=1, sticky=(W, E))

        self.excel = StringVar()
        excel_entry = ttk.Entry(mainframe, textvariable=self.excel)
        excel_entry.grid(column=2, row=2, sticky=(W, E))


        self.powerpoint = StringVar()
        powerpoint_entry = ttk.Entry(mainframe, textvariable=self.powerpoint)
        powerpoint_entry.grid(column=2, row=3, sticky=(W, E))


        ttk.Label(mainframe, text="Company Title").grid(column=1, row=1, sticky=W)
        ttk.Label(mainframe, text="Excel Location").grid(column=1, row=2, sticky=W)
        ttk.Label(mainframe, text="PowerPoint Location"). grid(column=1, row=3, sticky=W)
        
        ttk.Button(mainframe, text="Browse", command=self.browse_files_excel).grid(column=3, row=2, sticky=E)
        ttk.Button(mainframe, text="Browse", command=self.browse_files_powerpoint).grid(column=3, row=3, sticky=E)


        ttk.Button(mainframe, text="Run", command=self.run_program).grid(column=2, row=4)


    def browse_files_excel(self):
        filename = filedialog.askopenfilename(initialdir = "/", title = "Select an Excel", filetypes = [("Excel Files",".xlsx .xls")])
        self.excel.set(filename)

    def browse_files_powerpoint(self):
        filename = filedialog.askopenfilename(initialdir = "/", title = "Select a PowerPoint", filetypes = [("PowerPoint Files",".pptx .ppt")])
        self.powerpoint.set(filename)

    def excel_runner(self, loc):
        excel_wb = openpyxl.load_workbook(filename=loc)
        bw_util_sheet = excel_wb['BW Utilization']

        highly_utilized = []
        for row in bw_util_sheet.iter_rows(min_row=3, min_col=4, max_col=9, values_only=False):
            print(row[5])
            if row[5] != None and row[5].value >= .01:
                row[5].fill = PatternFill(fill_type='solid', fgColor='00FFFF00')
                highly_utilized.append(row[0].value)
        excel_wb.save(loc)
        return highly_utilized

    def powerpoint_runner(self, loc, circuits, entered_title):
        presentation = Presentation(loc)
        slide = presentation.slides[0]
        company_title = slide.shapes.title
        company_title.text = entered_title
        top = Inches(5.5)
        left = Inches(.1)
        current_slide = 1
        slide = presentation.slides[1]
        for circuit in circuits:
            circuit_label = slide.placeholders[1]
            circuit_label.text = circuit
            if current_slide < len(circuits):
                current_slide += 1
                slide = presentation.slides[current_slide]
        presentation.save(loc)


    def run_program(self):
        excel_loc = self.excel.get()
        powerpoint_loc = self.powerpoint.get()
        entered_title = self.title.get()

        circuits = self.excel_runner(excel_loc)
        print(circuits)

        self.powerpoint_runner(powerpoint_loc, circuits, entered_title)


root = Tk()
BandwidthPaster(root)
root.mainloop()

# TODO: Question: Can I get chart from API or do I have to Screenshot?

# TODO: Question: What does the security of the application have to look like for IT Requirements of the team?
            # Ideas - Users have to get individual API keys
                #   - Application works with Verizon SSO.

# TODO: Question: How do I package and bundle up the application into an executable